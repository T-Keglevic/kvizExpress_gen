"""
Quiz Slide Generator
Reads a text file with Q&A pairs and generates a PowerPoint presentation.
Questions and answers separated by blank lines.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from tqdm import tqdm
import sys


def parse_questions(filename):
    """
    Read text file and parse into list of (question, answer) tuples.
    Expected format:
        Question text here?
        Answer text here
        
        Next question?
        Next answer
    """
    with open(filename, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by double newlines (blank lines separate Q&A pairs)
    blocks = content.strip().split('\n\n')
    
    qa_pairs = []
    for block in blocks:
        lines = block.strip().split('\n')
        if len(lines) >= 2:
            question = lines[0].strip()
            answer = lines[1].strip()
            qa_pairs.append((question, answer))
        elif len(lines) == 1:
            # Question without answer - include it anyway with empty answer
            qa_pairs.append((lines[0].strip(), ""))
    
    return qa_pairs


def create_slide(prs, title_text, body_text, make_bold=False):
    """
    Create a slide with dark theme and centered title.
    
    Args:
        prs: Presentation object
        title_text: Text for top-centered title (question number)
        body_text: Main text content (question or answer)
        make_bold: Whether to make body text bold
    """
    # Use blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(32, 32, 32)  # Dark gray
    
    # Add title textbox (question number) - top centered
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title_text
    
    # Format title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.JUSTIFY
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)  # White
    
    # Add body textbox (question/answer content) - centered
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    
    body_box = slide.shapes.add_textbox(left, top, width, height)
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.text = body_text
    
    # Format body text
    body_para = body_frame.paragraphs[0]
    body_para.alignment = PP_ALIGN.JUSTIFY
    body_para.font.size = Pt(28)
    body_para.font.bold = make_bold
    body_para.font.color.rgb = RGBColor(255, 255, 255)  # White
    
    return slide


def generate_presentation(qa_pairs, output_filename):
    """
    Generate PowerPoint with question slides followed by answer slides.
    """
    prs = Presentation()
    
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    num_questions = len(qa_pairs)
    
    print(f"Generating {num_questions} question slides...")
    
    # First pass: Create all question slides
    for i, (question, answer) in tqdm(enumerate(qa_pairs, start=1), total=num_questions, desc="Questions"):
        title = f"{i}."
        create_slide(prs, title, question, make_bold=False)
    
    print(f"Generating {num_questions} answer slides...")
    
    # Second pass: Create all answer slides
    for i, (question, answer) in tqdm(enumerate(qa_pairs, start=1), total=num_questions, desc="Answers"):
        title = f"{i}."
        # Show question and answer together, answer in bold
        combined_text = f"{question}\n\n{answer}"
        
        # Create slide with question in normal, answer in bold
        # Note: python-pptx makes it tricky to mix bold/non-bold in one textbox
        # So we'll create two separate textboxes
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(32, 32, 32)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.JUSTIFY
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Question text (not bold)
        q_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        q_frame = q_box.text_frame
        q_frame.word_wrap = True
        q_frame.text = question
        q_para = q_frame.paragraphs[0]
        q_para.alignment = PP_ALIGN.JUSTIFY
        q_para.font.size = Pt(28)
        q_para.font.bold = False
        q_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Answer text (bold)
        a_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
        a_frame = a_box.text_frame
        a_frame.word_wrap = True
        a_frame.text = answer
        a_para = a_frame.paragraphs[0]
        a_para.alignment = PP_ALIGN.JUSTIFY
        a_para.font.size = Pt(28)
        a_para.font.bold = True
        a_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # Save presentation
    prs.save(output_filename)
    print(f"\n✓ Presentation saved as: {output_filename}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print(f"  - Questions: slides 1-{num_questions}")
    print(f"  - Answers: slides {num_questions + 1}-{len(prs.slides)}")


def main():
    """Main entry point for the script."""
    if len(sys.argv) != 2:
        print("Usage: python quiz_generator.py <questions_file.txt>")
        print("\nExpected format in text file:")
        print("  Question 1 text here?")
        print("  Answer 1 text here")
        print("  ")
        print("  Question 2 text here?")
        print("  Answer 2 text here")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = input_file.replace('.txt', '.pptx')
    
    # If input didn't have .txt extension, just append .pptx
    if output_file == input_file:
        output_file = input_file + '.pptx'
    
    print(f"Reading questions from: {input_file}")
    
    try:
        qa_pairs = parse_questions(input_file)
        
        if not qa_pairs:
            print("ERROR: No question/answer pairs found in file.")
            print("Check your file format - questions and answers should be separated by blank lines.")
            sys.exit(1)
        
        print(f"Found {len(qa_pairs)} question/answer pairs\n")
        
        generate_presentation(qa_pairs, output_file)
        
    except FileNotFoundError:
        print(f"ERROR: File '{input_file}' not found.")
        sys.exit(1)
    except Exception as e:
        print(f"ERROR: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()

## How to Use:

### 1. Create your questions file (`quiz.txt`):

#Question?
#Answer

#Who wrote Hamlet?
#William Shakespeare

### 2. Run the script:
# python quiz_generator.py quiz.txt